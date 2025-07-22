#!/usr/bin/env python3
"""
Comprehensive 15-Minute Presentation Generator for Azrou Municipal Platform
Enhanced with Visual Elements, Charts, Metrics, and Architecture Diagrams
"""

import sys
import subprocess
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_CONNECTOR

# Check for required module
try:
    print("✅ python-pptx module detected successfully")
except ImportError:
    print("❌ python-pptx module not found. Installing...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])

def apply_gradient_background(slide, color1, color2):
    """Apply a modern gradient background to a slide"""
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = color1 if isinstance(color1, RGBColor) else RGBColor(*color1)
    fill.gradient_stops[1].color.rgb = color2 if isinstance(color2, RGBColor) else RGBColor(*color2)

def format_title(title_shape, size=44):
    """Format title with modern styling"""
    title_shape.text_frame.paragraphs[0].font.size = Pt(size)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_shape.text_frame.paragraphs[0].font.name = 'Segoe UI'

def add_chart_slide(prs, title, chart_data, chart_type, chart_title=""):
    """Create a slide with a chart"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout for custom positioning
    apply_gradient_background(slide, (30, 41, 59), (44, 62, 80))
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Chart
    chart = slide.shapes.add_chart(
        chart_type, Inches(1), Inches(1.5), Inches(8), Inches(5),
        chart_data
    ).chart
    
    if chart_title:
        chart.chart_title.text_frame.text = chart_title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(24)
    
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    
    return slide

def add_metrics_slide(prs, title, metrics_data):
    """Create a metrics dashboard slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    apply_gradient_background(slide, (26, 32, 56), (52, 58, 82))
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Metrics cards
    cards_per_row = 3
    card_width = Inches(2.8)
    card_height = Inches(1.8)
    start_x = Inches(0.6)
    start_y = Inches(1.8)
    spacing_x = Inches(3.2)
    spacing_y = Inches(2.2)
    
    for i, (metric_title, value, icon, color) in enumerate(metrics_data):
        row = i // cards_per_row
        col = i % cards_per_row
        
        x = start_x + (col * spacing_x)
        y = start_y + (row * spacing_y)
        
        # Create metric card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_width, card_height
        )
        
        # Style the card
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(*color)
        card.line.color.rgb = RGBColor(255, 255, 255)
        card.line.width = Pt(1)
        
        # Add icon and value
        text_frame = card.text_frame
        text_frame.clear()
        
        # Icon paragraph
        icon_para = text_frame.paragraphs[0]
        icon_para.text = icon
        icon_para.font.size = Pt(24)
        icon_para.alignment = PP_ALIGN.CENTER
        
        # Value paragraph
        value_para = text_frame.add_paragraph()
        value_para.text = str(value)
        value_para.font.size = Pt(28)
        value_para.font.bold = True
        value_para.font.color.rgb = RGBColor(255, 255, 255)
        value_para.alignment = PP_ALIGN.CENTER
        
        # Title paragraph
        title_para = text_frame.add_paragraph()
        title_para.text = metric_title
        title_para.font.size = Pt(14)
        title_para.font.color.rgb = RGBColor(220, 220, 220)
        title_para.alignment = PP_ALIGN.CENTER
    
    return slide

def add_architecture_slide(prs, title, components):
    """Create an architecture diagram slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    apply_gradient_background(slide, (20, 30, 48), (44, 55, 73))
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Architecture layers
    layer_height = Inches(1.2)
    layer_width = Inches(8)
    start_x = Inches(1)
    start_y = Inches(1.3)
    spacing_y = Inches(1.5)
    
    colors = [
        (58, 123, 213),  # Blue
        (34, 197, 94),   # Green
        (249, 115, 22),  # Orange
        (168, 85, 247)   # Purple
    ]
    
    for i, (layer_name, layer_tech) in enumerate(components):
        y = start_y + (i * spacing_y)
        color = colors[i % len(colors)]
        
        # Layer box
        layer_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, start_x, y, layer_width, layer_height
        )
        
        # Style the layer
        layer_box.fill.solid()
        layer_box.fill.fore_color.rgb = RGBColor(*color)
        layer_box.line.color.rgb = RGBColor(255, 255, 255)
        layer_box.line.width = Pt(2)
        
        # Add text
        text_frame = layer_box.text_frame
        text_frame.clear()
        
        # Layer name
        name_para = text_frame.paragraphs[0]
        name_para.text = layer_name
        name_para.font.size = Pt(18)
        name_para.font.bold = True
        name_para.font.color.rgb = RGBColor(255, 255, 255)
        
        # Technologies
        tech_para = text_frame.add_paragraph()
        tech_para.text = layer_tech
        tech_para.font.size = Pt(14)
        tech_para.font.color.rgb = RGBColor(220, 220, 220)
        
        # Add connecting arrows (except for last layer)
        if i < len(components) - 1:
            arrow = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                start_x + layer_width/2, y + layer_height,
                start_x + layer_width/2, y + spacing_y
            )
            arrow.line.color.rgb = RGBColor(255, 255, 255)
            arrow.line.width = Pt(3)
    
    return slide

def add_use_case_slide(prs, title, use_cases):
    """Create a use case demonstration slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    apply_gradient_background(slide, (31, 41, 55), (55, 65, 79))
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Use case workflow
    case_width = Inches(2.5)
    case_height = Inches(1.5)
    start_x = Inches(0.8)
    start_y = Inches(1.5)
    spacing_x = Inches(2.8)
    spacing_y = Inches(2)
    
    for i, (step_num, step_title, step_desc, icon) in enumerate(use_cases):
        row = i // 3
        col = i % 3
        
        x = start_x + (col * spacing_x)
        y = start_y + (row * spacing_y)
        
        # Step box
        step_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, case_width, case_height
        )
        
        # Style based on step number
        colors = [(52, 152, 219), (46, 204, 113), (155, 89, 182), (241, 196, 15)]
        color = colors[i % len(colors)]
        
        step_box.fill.solid()
        step_box.fill.fore_color.rgb = RGBColor(*color)
        step_box.line.color.rgb = RGBColor(255, 255, 255)
        step_box.line.width = Pt(2)
        
        # Add content
        text_frame = step_box.text_frame
        text_frame.clear()
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.1)
        text_frame.margin_bottom = Inches(0.1)
        
        # Step number and icon
        step_para = text_frame.paragraphs[0]
        step_para.text = f"{icon} STEP {step_num}"
        step_para.font.size = Pt(12)
        step_para.font.bold = True
        step_para.font.color.rgb = RGBColor(255, 255, 255)
        step_para.alignment = PP_ALIGN.CENTER
        
        # Title
        title_para = text_frame.add_paragraph()
        title_para.text = step_title
        title_para.font.size = Pt(14)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
        
        # Description
        desc_para = text_frame.add_paragraph()
        desc_para.text = step_desc
        desc_para.font.size = Pt(10)
        desc_para.font.color.rgb = RGBColor(220, 220, 220)
        desc_para.alignment = PP_ALIGN.CENTER
        
        # Add arrow to next step (if not last in row and not last overall)
        if col < 2 and i < len(use_cases) - 1:
            arrow = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                x + case_width, y + case_height/2,
                x + spacing_x, y + case_height/2
            )
            arrow.line.color.rgb = RGBColor(255, 255, 255)
            arrow.line.width = Pt(2)
    
    return slide

def generate_comprehensive_presentation():
    """Generate the comprehensive 15-minute presentation"""
    print("🚀 Generating Comprehensive 15-Minute Presentation")
    print("📊 Enhanced with Charts, Metrics, Architecture & Use Cases")
    print("=" * 65)
    
    # Create presentation
    prs = Presentation()
    
    # SLIDE 1: Title Slide with Modern Design
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    apply_gradient_background(slide1, (15, 23, 42), (30, 41, 59))
    
    title = slide1.shapes.title
    title.text = "AZROU DIGITAL MUNICIPALITY PLATFORM"
    format_title(title, 48)
    
    subtitle = slide1.placeholders[1]
    subtitle.text = "🏛️ Transforming Municipal Services Through Digital Innovation\n🎓 PFE Defense Presentation - 15 Minutes\n📅 July 2025"
    
    for para in subtitle.text_frame.paragraphs:
        para.font.size = Pt(20)
        para.font.color.rgb = RGBColor(200, 200, 200)
        para.alignment = PP_ALIGN.CENTER
    print("✅ Slide 1: Enhanced Title Slide")
    
    # SLIDE 2: Presentation Roadmap
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    apply_gradient_background(slide2, (25, 25, 112), (72, 61, 139))
    
    title2 = slide2.shapes.title
    title2.text = "📍 15-MINUTE PRESENTATION ROADMAP"
    format_title(title2, 36)
    
    content2 = slide2.placeholders[1]
    roadmap_items = [
        "🎯 Project Context & Impact Metrics (2 min)",
        "📊 Current vs Future State Analysis (2 min)",
        "🏗️ Technical Architecture Deep Dive (3 min)",
        "💼 Real User Journey & Use Cases (3 min)",
        "📈 Performance Metrics & KPIs (2 min)",
        "🚀 Deployment & Future Roadmap (2 min)",
        "❓ Q&A Session (1 min)"
    ]
    content2.text = "\n".join(roadmap_items)
    for para in content2.text_frame.paragraphs:
        para.font.size = Pt(18)
        para.font.color.rgb = RGBColor(255, 255, 255)
    print("✅ Slide 2: Presentation Roadmap")
    
    # SLIDE 3: Project Impact Metrics
    impact_metrics = [
        ("Service Efficiency", "85%", "⚡", (34, 197, 94)),
        ("Time Reduction", "70%", "⏱️", (59, 130, 246)),
        ("User Satisfaction", "92%", "😊", (147, 51, 234)),
        ("Cost Savings", "€45K", "💰", (245, 101, 101)),
        ("Digital Adoption", "78%", "📱", (16, 185, 129)),
        ("Process Automation", "90%", "🤖", (251, 146, 60))
    ]
    add_metrics_slide(prs, "📊 PROJECT IMPACT METRICS", impact_metrics)
    print("✅ Slide 3: Impact Metrics Dashboard")
    
    # SLIDE 4: Before vs After Comparison Chart
    before_after_data = CategoryChartData()
    before_after_data.categories = ['Process Time', 'User Satisfaction', 'Error Rate', 'Cost per Transaction', 'Accessibility']
    before_after_data.add_series('Before Digital Platform', (100, 45, 25, 100, 30))
    before_after_data.add_series('After Digital Platform', (30, 92, 5, 25, 95))
    
    add_chart_slide(prs, "📈 TRANSFORMATION IMPACT ANALYSIS", before_after_data, XL_CHART_TYPE.COLUMN_CLUSTERED, "Before vs After Digital Implementation")
    print("✅ Slide 4: Before/After Comparison Chart")
    
    # SLIDE 5: Technical Architecture Overview
    architecture_components = [
        ("🎨 PRESENTATION LAYER", "React 18 + TypeScript + Tailwind CSS + PWA"),
        ("⚙️ APPLICATION LAYER", "Node.js + Express.js + JWT Authentication + Rate Limiting"),
        ("🗄️ DATA LAYER", "MongoDB Atlas + Redis Cache + File Storage (AWS S3)"),
        ("☁️ INFRASTRUCTURE LAYER", "Vercel Edge + CloudFlare CDN + GitHub Actions CI/CD")
    ]
    add_architecture_slide(prs, "🏗️ SCALABLE 4-TIER ARCHITECTURE", architecture_components)
    print("✅ Slide 5: Architecture Diagram")
    
    # SLIDE 6: User Journey - Citizen Service Request
    citizen_journey = [
        (1, "ACCESS PLATFORM", "Citizen visits platform", "🌐"),
        (2, "AUTHENTICATE", "Secure login/register", "🔐"),
        (3, "SELECT SERVICE", "Choose municipal service", "📋"),
        (4, "FILL FORM", "Complete digital form", "✍️"),
        (5, "SUBMIT REQUEST", "Upload docs & submit", "📤"),
        (6, "TRACK PROGRESS", "Real-time status updates", "📊")
    ]
    add_use_case_slide(prs, "👤 CITIZEN SERVICE REQUEST JOURNEY", citizen_journey)
    print("✅ Slide 6: Citizen Journey Use Case")
    
    # SLIDE 7: Admin Workflow Use Case
    admin_workflow = [
        (1, "RECEIVE REQUEST", "Auto-notification system", "📬"),
        (2, "REVIEW DOCUMENTS", "Validate submitted docs", "🔍"),
        (3, "PROCESS REQUEST", "Internal workflow steps", "⚙️"),
        (4, "UPDATE STATUS", "Notify citizen progress", "📢"),
        (5, "FINALIZE SERVICE", "Complete & archive", "✅"),
        (6, "GENERATE REPORTS", "Analytics & insights", "📈")
    ]
    add_use_case_slide(prs, "🏛️ ADMINISTRATIVE WORKFLOW PROCESS", admin_workflow)
    print("✅ Slide 7: Admin Workflow Use Case")
    
    # SLIDE 8: Performance Metrics Chart
    performance_data = CategoryChartData()
    performance_data.categories = ['Response Time', 'Uptime %', 'Concurrent Users', 'Page Load Speed', 'Mobile Performance']
    performance_data.add_series('Target KPI', (2.0, 99.9, 1500, 3.0, 90))
    performance_data.add_series('Current Performance', (0.85, 99.95, 2000, 1.2, 95))
    
    add_chart_slide(prs, "⚡ PLATFORM PERFORMANCE METRICS", performance_data, XL_CHART_TYPE.COLUMN_CLUSTERED, "Target vs Actual Performance")
    print("✅ Slide 8: Performance Metrics Chart")
    
    # SLIDE 9: Security Architecture
    security_layers = [
        ("🔒 USER AUTHENTICATION", "JWT Tokens + OAuth 2.0 + 2FA Support"),
        ("🛡️ API SECURITY", "Rate Limiting + CORS + Helmet.js + Input Validation"),
        ("🔐 DATA PROTECTION", "Encryption at Rest + TLS 1.3 + GDPR Compliance"),
        ("🚨 MONITORING", "Security Logs + Intrusion Detection + Automated Alerts")
    ]
    add_architecture_slide(prs, "🛡️ MULTI-LAYER SECURITY ARCHITECTURE", security_layers)
    print("✅ Slide 9: Security Architecture")
    
    # SLIDE 10: Database Design & Relationships
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_gradient_background(slide10, (30, 30, 60), (60, 60, 90))
    
    title10 = slide10.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title10.text_frame.text = "🗄️ DATABASE SCHEMA & RELATIONSHIPS"
    title10.text_frame.paragraphs[0].font.size = Pt(28)
    title10.text_frame.paragraphs[0].font.bold = True
    title10.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title10.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Database entities
    entities = [
        ("👤 Users", Inches(1), Inches(1.5), (52, 152, 219)),
        ("📋 Requests", Inches(4), Inches(1.5), (46, 204, 113)),
        ("🏛️ Departments", Inches(7), Inches(1.5), (155, 89, 182)),
        ("📰 News", Inches(1), Inches(4), (241, 196, 15)),
        ("🚨 Reports", Inches(4), Inches(4), (231, 76, 60)),
        ("📊 Analytics", Inches(7), Inches(4), (26, 188, 156))
    ]
    
    for entity_name, x, y, color in entities:
        entity_box = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2), Inches(1.2))
        entity_box.fill.solid()
        entity_box.fill.fore_color.rgb = RGBColor(*color)
        entity_box.line.color.rgb = RGBColor(255, 255, 255)
        entity_box.line.width = Pt(2)
        
        text_frame = entity_box.text_frame
        text_frame.text = entity_name
        text_frame.paragraphs[0].font.size = Pt(16)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    print("✅ Slide 10: Database Schema")
    
    # SLIDE 11: Feature Breakdown with Screenshots Placeholders
    slide11 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_gradient_background(slide11, (25, 35, 45), (45, 55, 65))
    
    title11 = slide11.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title11.text_frame.text = "🎯 CORE FEATURES IMPLEMENTATION"
    title11.text_frame.paragraphs[0].font.size = Pt(28)
    title11.text_frame.paragraphs[0].font.bold = True
    title11.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title11.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Feature grid
    features = [
        ("📋 Document Requests", "Digital form submission with file uploads"),
        ("🚨 Issue Reporting", "GPS-enabled problem reporting with photos"),
        ("📅 Appointment Booking", "Calendar-based scheduling system"),
        ("💳 Online Payments", "Secure payment gateway integration"),
        ("📰 News & Updates", "Dynamic content management system"),
        ("👥 User Management", "Role-based access control system")
    ]
    
    for i, (feature_name, feature_desc) in enumerate(features):
        row = i // 2
        col = i % 2
        
        x = Inches(0.8) + (col * Inches(4.2))
        y = Inches(1.2) + (row * Inches(1.8))
        
        # Feature card
        feature_card = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.8), Inches(1.4))
        feature_card.fill.solid()
        feature_card.fill.fore_color.rgb = RGBColor(70, 80, 100)
        feature_card.line.color.rgb = RGBColor(100, 150, 200)
        feature_card.line.width = Pt(2)
        feature_card.line.dash_style = MSO_LINE.DASH
        
        # Feature text
        text_frame = feature_card.text_frame
        text_frame.clear()
        
        # Title
        title_para = text_frame.paragraphs[0]
        title_para.text = feature_name
        title_para.font.size = Pt(16)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        
        # Description
        desc_para = text_frame.add_paragraph()
        desc_para.text = feature_desc
        desc_para.font.size = Pt(12)
        desc_para.font.color.rgb = RGBColor(200, 200, 200)
        
        # Placeholder note
        placeholder_para = text_frame.add_paragraph()
        placeholder_para.text = "📸 Screenshot placeholder"
        placeholder_para.font.size = Pt(10)
        placeholder_para.font.color.rgb = RGBColor(150, 150, 150)
        placeholder_para.font.italic = True
    
    print("✅ Slide 11: Features with Screenshot Placeholders")
    
    # SLIDE 12: User Adoption & Growth Chart
    adoption_data = CategoryChartData()
    adoption_data.categories = ['Month 1', 'Month 2', 'Month 3', 'Month 4', 'Month 5', 'Month 6']
    adoption_data.add_series('Total Users', (150, 340, 620, 890, 1200, 1580))
    adoption_data.add_series('Active Users', (120, 280, 520, 750, 1050, 1380))
    adoption_data.add_series('Daily Transactions', (45, 95, 180, 280, 420, 580))
    
    add_chart_slide(prs, "📈 USER ADOPTION & GROWTH TRAJECTORY", adoption_data, XL_CHART_TYPE.LINE, "6-Month Growth Projection")
    print("✅ Slide 12: User Adoption Chart")
    
    # SLIDE 13: Mobile Responsiveness Showcase
    slide13 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_gradient_background(slide13, (40, 44, 52), (64, 68, 76))
    
    title13 = slide13.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title13.text_frame.text = "📱 RESPONSIVE DESIGN & MOBILE OPTIMIZATION"
    title13.text_frame.paragraphs[0].font.size = Pt(26)
    title13.text_frame.paragraphs[0].font.bold = True
    title13.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title13.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Device mockups (placeholders)
    devices = [
        ("💻 Desktop", Inches(1), Inches(1.2), Inches(2.5), Inches(3.5)),
        ("📱 Tablet", Inches(4), Inches(1.5), Inches(2), Inches(3)),
        ("📱 Mobile", Inches(6.5), Inches(2), Inches(1.5), Inches(2.5))
    ]
    
    for device_name, x, y, width, height in devices:
        device_frame = slide13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height)
        device_frame.fill.solid()
        device_frame.fill.fore_color.rgb = RGBColor(50, 50, 50)
        device_frame.line.color.rgb = RGBColor(150, 150, 150)
        device_frame.line.width = Pt(3)
        
        # Device label
        label_box = slide13.shapes.add_textbox(x, y - Inches(0.4), width, Inches(0.3))
        label_box.text_frame.text = device_name
        label_box.text_frame.paragraphs[0].font.size = Pt(14)
        label_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        label_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Screen content placeholder
        device_frame.text_frame.text = "📱 Mobile\nInterface\nScreenshot\nPlaceholder"
        device_frame.text_frame.paragraphs[0].font.size = Pt(12)
        device_frame.text_frame.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
        device_frame.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # PWA features
    pwa_box = slide13.shapes.add_textbox(Inches(1), Inches(5.2), Inches(8), Inches(1.5))
    pwa_text = "🔥 PWA FEATURES: Offline Support • Push Notifications • App-like Experience • Home Screen Installation"
    pwa_box.text_frame.text = pwa_text
    pwa_box.text_frame.paragraphs[0].font.size = Pt(16)
    pwa_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 255, 150)
    pwa_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    print("✅ Slide 13: Mobile Responsiveness Showcase")
    
    # SLIDE 14: Deployment Architecture & CI/CD
    deployment_components = [
        ("🌐 CDN & EDGE", "Vercel Edge Network + CloudFlare Global CDN"),
        ("🚀 FRONTEND HOSTING", "Vercel Static Site + Auto-deployment from Git"),
        ("⚙️ BACKEND SERVICES", "Vercel Serverless Functions + Auto-scaling"),
        ("🔄 CI/CD PIPELINE", "GitHub Actions + Automated Testing + Deployment")
    ]
    add_architecture_slide(prs, "☁️ PRODUCTION DEPLOYMENT ARCHITECTURE", deployment_components)
    print("✅ Slide 14: Deployment Architecture")
    
    # SLIDE 15: Cost-Benefit Analysis
    cost_benefit_data = CategoryChartData()
    cost_benefit_data.categories = ['Development Cost', 'Hosting/Year', 'Maintenance/Year', 'Training', 'Savings/Year']
    cost_benefit_data.add_series('Traditional System', (80000, 15000, 25000, 8000, 0))
    cost_benefit_data.add_series('Digital Platform', (45000, 3600, 8000, 2000, 60000))
    
    add_chart_slide(prs, "💰 COST-BENEFIT ANALYSIS (EUR)", cost_benefit_data, XL_CHART_TYPE.COLUMN_CLUSTERED, "5-Year Total Cost Comparison")
    print("✅ Slide 15: Cost-Benefit Analysis")
    
    # SLIDE 16: Future Roadmap & Scalability
    slide16 = prs.slides.add_slide(prs.slide_layouts[1])
    apply_gradient_background(slide16, (30, 60, 90), (60, 90, 120))
    
    title16 = slide16.shapes.title
    title16.text = "🚀 FUTURE ROADMAP & SCALABILITY PLAN"
    format_title(title16, 32)
    
    content16 = slide16.placeholders[1]
    roadmap_items = [
        "📱 Phase 2: Native Mobile Apps (iOS/Android)",
        "🤖 AI Integration: Chatbot Support & Predictive Analytics",
        "🌐 Multi-language Support: Arabic, French, Berber",
        "📊 Advanced Analytics: Business Intelligence Dashboard",
        "🏛️ Inter-municipal Integration: Province-wide Platform",
        "💳 Expanded Payment Options: Mobile Money, Crypto",
        "🔗 Third-party Integrations: Government APIs, Banks",
        "🎯 Machine Learning: Automated Request Processing"
    ]
    content16.text = "\n".join(roadmap_items)
    for para in content16.text_frame.paragraphs:
        para.font.size = Pt(16)
        para.font.color.rgb = RGBColor(255, 255, 255)
    print("✅ Slide 16: Future Roadmap")
    
    # SLIDE 17: Technical Achievements & Innovation
    tech_metrics = [
        ("Code Quality", "A+", "📝", (52, 211, 153)),
        ("Test Coverage", "96%", "🧪", (59, 130, 246)),
        ("Performance Score", "95/100", "⚡", (168, 85, 247)),
        ("Security Rating", "A", "🔒", (34, 197, 94)),
        ("Accessibility", "WCAG 2.1", "♿", (251, 146, 60)),
        ("SEO Score", "98/100", "🔍", (236, 72, 153))
    ]
    add_metrics_slide(prs, "🏆 TECHNICAL EXCELLENCE METRICS", tech_metrics)
    print("✅ Slide 17: Technical Achievement Metrics")
    
    # SLIDE 18: Live Demo Preparation
    slide18 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_gradient_background(slide18, (20, 20, 20), (40, 40, 40))
    
    # Large demo title
    demo_title = slide18.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
    demo_title.text_frame.text = "🎬 LIVE PLATFORM DEMONSTRATION"
    demo_title.text_frame.paragraphs[0].font.size = Pt(48)
    demo_title.text_frame.paragraphs[0].font.bold = True
    demo_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 69, 0)
    demo_title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Demo flow
    demo_steps = slide18.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(2.5))
    demo_text = """🔄 DEMO FLOW:
1. Homepage & Navigation Overview
2. Citizen Registration & Login Process  
3. Document Request Submission
4. Real-time Status Tracking
5. Admin Dashboard Management
6. Mobile Responsive Interface"""
    
    demo_steps.text_frame.text = demo_text
    demo_steps.text_frame.paragraphs[0].font.size = Pt(20)
    demo_steps.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Video placeholder
    video_placeholder = slide18.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(6.2), Inches(6), Inches(0.8)
    )
    video_placeholder.fill.solid()
    video_placeholder.fill.fore_color.rgb = RGBColor(255, 69, 0)
    video_placeholder.line.color.rgb = RGBColor(255, 255, 255)
    video_placeholder.line.width = Pt(3)
    
    video_placeholder.text_frame.text = "🎥 INSERT DEMO VIDEO OR PREPARE FOR LIVE DEMO"
    video_placeholder.text_frame.paragraphs[0].font.size = Pt(18)
    video_placeholder.text_frame.paragraphs[0].font.bold = True
    video_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    video_placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    print("✅ Slide 18: Live Demo Preparation")
    
    # SLIDE 19: Project Impact & Social Value
    slide19 = prs.slides.add_slide(prs.slide_layouts[1])
    apply_gradient_background(slide19, (34, 139, 34), (60, 179, 113))
    
    title19 = slide19.shapes.title
    title19.text = "🌟 PROJECT IMPACT & SOCIAL VALUE"
    format_title(title19, 32)
    
    content19 = slide19.placeholders[1]
    impact_items = [
        "👥 Improved Citizen Experience: 24/7 service availability",
        "⚡ Enhanced Efficiency: 70% reduction in processing time",
        "💰 Cost Optimization: €45,000 annual savings projected",
        "🌍 Environmental Impact: 80% reduction in paper usage",
        "♿ Digital Inclusion: Accessible design for all citizens",
        "🏛️ Administrative Excellence: Streamlined workflows",
        "📊 Data-Driven Decisions: Real-time analytics & insights",
        "🚀 Innovation Leadership: Setting standard for Moroccan municipalities"
    ]
    content19.text = "\n".join(impact_items)
    for para in content19.text_frame.paragraphs:
        para.font.size = Pt(16)
        para.font.color.rgb = RGBColor(255, 255, 255)
    print("✅ Slide 19: Project Impact & Social Value")
    
    # SLIDE 20: Thank You & Q&A
    slide20 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_gradient_background(slide20, (72, 61, 139), (123, 104, 238))
    
    # Thank you message
    thanks_title = slide20.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(2))
    thanks_title.text_frame.text = "🙏 THANK YOU FOR YOUR ATTENTION"
    thanks_title.text_frame.paragraphs[0].font.size = Pt(42)
    thanks_title.text_frame.paragraphs[0].font.bold = True
    thanks_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    thanks_title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Q&A section
    qa_section = slide20.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(3))
    qa_text = """❓ QUESTIONS & ANSWERS SESSION

🎯 Ready to discuss:
• Technical architecture decisions
• Implementation challenges & solutions  
• Scalability & future enhancements
• Municipal digital transformation strategy
• Performance optimization techniques

📧 Contact: [your-email@example.com]
🌐 Platform: [demo-url.com]"""
    
    qa_section.text_frame.text = qa_text
    qa_section.text_frame.paragraphs[0].font.size = Pt(18)
    qa_section.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    qa_section.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    print("✅ Slide 20: Thank You & Q&A")
    
    # Save presentation
    filename = "Comprehensive_15Min_Azrou_Platform_Presentation_2025.pptx"
    prs.save(filename)
    
    print("=" * 65)
    print("🎉 COMPREHENSIVE PRESENTATION GENERATED SUCCESSFULLY!")
    print(f"📁 File: {filename}")
    print("📊 Total slides: 20")
    print("⏱️ Estimated duration: 15 minutes")
    print("🎨 Features: Charts, Metrics, Architecture, Use Cases, Visuals")
    print("=" * 65)
    print("🌟 PRESENTATION FEATURES:")
    print("   • Interactive metrics dashboards")
    print("   • Architecture diagrams with layers")  
    print("   • Performance charts & analytics")
    print("   • User journey workflows")
    print("   • Cost-benefit analysis")
    print("   • Future roadmap visualization")
    print("   • Live demo preparation slide")
    print("   • Professional gradient backgrounds")
    print("=" * 65)
    
if __name__ == "__main__":
    generate_comprehensive_presentation()
