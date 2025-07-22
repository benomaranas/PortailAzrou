#!/usr/bin/env python3
"""
Extended PowerPoint Presentation Generator for Final Year Project (30+ slides)
Municipal Platform Azrou - Anas BENOMAR
Complete Version for University Defense
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.dml import MSO_THEME_COLOR
    print("✅ python-pptx module detected")
except ImportError:
    print("❌ python-pptx module not installed")
    print("🔧 Required installation: pip install python-pptx")
    exit(1)

def add_styled_title(slide, title_text, color_rgb=(0, 51, 102)):
    """Add a styled title to a slide"""
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(*color_rgb)
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return title

def add_bullet_points(content_placeholder, points_list):
    """Add bullet points to a placeholder"""
    content_placeholder.text = points_list[0]
    for point in points_list[1:]:
        p = content_placeholder.text_frame.add_paragraph()
        p.text = point
        p.level = 0
    
    # Bullet point styling
    for paragraph in content_placeholder.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.space_after = Pt(6)

def add_section_header(prs, section_title, section_subtitle="", color_rgb=(34, 139, 34)):
    """Add a section slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[2])  # Section header layout
    title = slide.shapes.title
    title.text = section_title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(*color_rgb)
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    if section_subtitle and len(slide.placeholders) > 1:
        subtitle = slide.placeholders[1]
        subtitle.text = section_subtitle
        subtitle.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def create_comprehensive_presentation():
    """Create a comprehensive PowerPoint presentation with 30+ slides"""
    
    # Create a new presentation
    prs = Presentation()
    
    print("🚀 Generating extended presentation with 30+ slides...")
    print("=" * 60)
    
    # ============== SLIDE 1: TITLE PAGE ==============
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title1 = slide1.shapes.title
    subtitle1 = slide1.placeholders[1]
    
    title1.text = "Development of an Integrated Municipal Digital Platform"
    subtitle1.text = """Municipality of Azrou - Public Services Digitalization

Presented by: Anas BENOMAR
Academic Supervisor: Mr. MEKOUAR Othmane
Professional Supervisor: IT Services Municipality of Azrou

Higher School of Engineering in Applied Sciences (ESISA)
Final Year Project - Academic Year 2024/2025"""
    
    # Style title
    title1.text_frame.paragraphs[0].font.size = Pt(32)
    title1.text_frame.paragraphs[0].font.bold = True
    title1.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Style subtitle
    for p in subtitle1.text_frame.paragraphs:
        p.font.size = Pt(16)
        p.alignment = PP_ALIGN.CENTER
    
    print("✅ Slide 1: Title page")
    
    # ============== SLIDE 2: PRESENTATION OUTLINE ==============
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    add_styled_title(slide2, "Presentation Outline")
    
    content2 = slide2.placeholders[1]
    plan_points = [
        "I. Context and Problem Statement",
        "II. Needs Analysis and State of the Art",
        "III. System Design and Architecture",
        "IV. Development and Implementation",
        "V. Platform Demonstration",
        "VI. Testing and Validation",
        "VII. Deployment and Maintenance",
        "VIII. Impact and Future Perspectives",
        "IX. Conclusion and Questions"
    ]
    add_bullet_points(content2, plan_points)
    print("✅ Slide 2: General outline")
    
    # ============== SECTION I: CONTEXT ==============
    add_section_header(prs, "I. CONTEXT AND PROBLEM STATEMENT", "Municipal Services Digitalization")
    print("✅ Section I: Header")
    
    # SLIDE 3: Presentation of Azrou Municipality
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    add_styled_title(slide3, "Municipality of Azrou - Overview")
    
    content3 = slide3.placeholders[1]
    azrou_points = [
        "🏛️ Urban municipality in Ifrane province",
        "👥 Population: ~60,000 inhabitants",
        "🌍 Located in the heart of Middle Atlas",
        "🏗️ Important regional economic center",
        "📊 Diverse and complex municipal services",
        "🔄 Need for digital modernization",
        "🎯 Commitment to improve administrative efficiency"
    ]
    add_bullet_points(content3, azrou_points)
    print("✅ Slide 3: Azrou presentation")
    
    # SLIDE 4: Identified Problems
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    add_styled_title(slide4, "Current Challenges")
    
    content4 = slide4.placeholders[1]
    problems_points = [
        "⏰ Long and complex administrative processes",
        "📋 Paper-based request and file management",
        "🚶‍♂️ Mandatory travel for citizens",
        "📞 Lack of modern communication channels",
        "🔍 Difficulties in tracking requests",
        "📊 Absence of analysis and reporting tools",
        "💼 Heavy administrative burden",
        "🌐 Lack of strong digital presence"
    ]
    add_bullet_points(content4, problems_points)
    print("✅ Slide 4: Problems")
    
    # SLIDE 5: Project Objectives
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    add_styled_title(slide5, "Project Objectives")
    
    content5 = slide5.placeholders[1]
    objectives_points = [
        "🎯 Digitalize essential municipal services",
        "⚡ Simplify administrative procedures",
        "🏠 Enable online requests from home",
        "📱 Create a modern and accessible interface",
        "🔧 Facilitate internal service management",
        "📈 Improve operational efficiency",
        "🤝 Strengthen citizen-administration relationship",
        "🌱 Lay the foundation for a Smart City"
    ]
    add_bullet_points(content5, objectives_points)
    print("✅ Slide 5: Objectives")
    
    # ============== SECTION II: ANALYSIS ==============
    add_section_header(prs, "II. ANALYSIS AND STATE OF THE ART", "Needs Study and Existing Solutions")
    print("✅ Section II: Header")
    
    # SLIDE 6: Analysis Methodology
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    add_styled_title(slide6, "Needs Analysis Methodology")
    
    content6 = slide6.placeholders[1]
    methodology_points = [
        "🗣️ Interviews with municipal agents",
        "👥 Citizen surveys",
        "🔍 Observation of existing processes",
        "📊 Analysis of request volumes",
        "⏱️ Measurement of processing times",
        "💻 Audit of current IT systems",
        "🌐 Benchmarking of existing solutions",
        "📋 Workflow documentation"
    ]
    add_bullet_points(content6, methodology_points)
    print("✅ Slide 6: Methodology")
    
    # SLIDE 7: Identified Functional Requirements
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    add_styled_title(slide7, "Functional Requirements")
    
    content7 = slide7.placeholders[1]
    functional_needs = [
        "📝 Administrative document requests",
        "🏗️ Construction permit applications",
        "💡 Urban problem reporting",
        "📰 Municipal news publication",
        "📅 Event and manifestation management",
        "👤 Citizen account management",
        "🔐 Secure authentication system",
        "📊 Administrative dashboards"
    ]
    add_bullet_points(content7, functional_needs)
    print("✅ Slide 7: Functional requirements")
    
    # SLIDE 8: Non-Functional Requirements
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    add_styled_title(slide8, "Non-Functional Requirements")
    
    content8 = slide8.placeholders[1]
    nonfunctional_needs = [
        "🚀 Performance: Response time < 3 seconds",
        "👥 Scalability: Support 1000+ users",
        "🔒 Security: Encryption and authentication",
        "📱 Responsiveness: All device compatibility",
        "♿ Accessibility: Web standards compliance",
        "🌐 Availability: 99.9% uptime",
        "💾 Backup: Data protection",
        "🔧 Maintainability: Modular architecture"
    ]
    add_bullet_points(content8, nonfunctional_needs)
    print("✅ Slide 8: Non-functional requirements")
    
    # SLIDE 9: State of the Art - Competing Solutions
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    add_styled_title(slide9, "State of the Art - Municipal Platforms")
    
    content9 = slide9.placeholders[1]
    competitors_points = [
        "🇫🇷 Démarches Simplifiées (France)",
        "🇪🇸 Carpeta Ciudadana (Spain)",
        "🇨🇦 Mon Dossier Citoyen (Canada)",
        "🇲🇦 Chikaya.ma (Morocco) - Reports",
        "🏛️ Various GovTech Solutions",
        "💼 Existing proprietary systems",
        "🔍 Analysis of strengths and weaknesses",
        "🎯 Identification of improvement areas"
    ]
    add_bullet_points(content9, competitors_points)
    print("✅ Slide 9: State of the art")
    
    # SLIDE 10: Justified Technology Choices
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    add_styled_title(slide10, "Technology Choices")
    
    content10 = slide10.placeholders[1]
    tech_choices = [
        "⚛️ Frontend: React.js + Tailwind CSS",
        "🚀 Backend: Node.js + Express.js",
        "🗄️ Database: MongoDB",
        "☁️ Cloud: AWS/Vercel services",
        "🔐 Authentication: JWT + bcrypt",
        "📧 Notifications: Integrated email service",
        "📱 Design: Responsive and mobile-first",
        "🔧 Tools: Git, VSCode, Postman"
    ]
    add_bullet_points(content10, tech_choices)
    print("✅ Slide 10: Technology choices")
    
    # ============== SECTION III: CONCEPTION ==============
    add_section_header(prs, "III. CONCEPTION SYSTÈME", "Architecture et Design de la Solution")
    print("✅ Section III: Header")
    
    # SLIDE 11: Architecture Globale
    slide11 = prs.slides.add_slide(prs.slide_layouts[1])
    add_styled_title(slide11, "Architecture Globale du Système")
    
    content11 = slide11.placeholders[1]
    architecture_points = [
        "🌐 Architecture 3-tiers moderne",
        "👤 Couche Présentation: Interface React",
        "⚙️ Couche Métier: API REST Node.js",
        "🗄️ Couche Données: MongoDB Atlas",
        "🔄 Communication: API RESTful + JSON",
        "🔐 Sécurité: Middleware d'authentification",
        "📱 Design Pattern: MVC côté backend",
        "🎨 State Management: Context API React"
    ]
    add_bullet_points(content11, architecture_points)
    print("✅ Slide 11: Architecture globale")
    
    # SLIDE 12: Modélisation Base de Données
    slide12 = prs.slides.add_slide(prs.slide_layouts[1])
    add_styled_title(slide12, "Modélisation de la Base de Données")
    
    content12 = slide12.placeholders[1]
    database_points = [
        "👤 Collection Users: Profils citoyens",
        "📄 Collection Documents: Demandes administratives",
        "🚨 Collection Reports: Signalements citoyens",
        "📰 Collection News: Actualités municipales",
        "📅 Collection Events: Événements publics",
        "🔗 Relations: Références croisées optimisées",
        "📊 Indexation: Performance des requêtes",
        "🔄 Schémas flexibles: Évolutivité MongoDB"
    ]
    add_bullet_points(content12, database_points)
    print("✅ Slide 12: Base de données")
    
    # SLIDE 13: Design UX/UI
    slide13 = prs.slides.add_slide(prs.slide_layouts[1])
    add_styled_title(slide13, "Conception UX/UI")
    
    content13 = slide13.placeholders[1]
    design_points = [
        "🎨 Design System: Couleurs municipales",
        "📱 Mobile-First: Responsive design",
        "🧭 Navigation intuitive: Menu clair",
        "♿ Accessibilité: WCAG 2.1 compliance",
        "🔍 Recherche: Moteur intégré efficace",
        "📋 Formulaires: Validation en temps réel",
        "🔔 Notifications: Feedback utilisateur",
        "⚡ Performance: Temps de chargement optimisés"
    ]
    add_bullet_points(content13, design_points)
    print("✅ Slide 13: Design UX/UI")
    
    # SLIDE 14: Diagrammes de Cas d'Usage
    slide14 = prs.slides.add_slide(prs.slide_layouts[1])
    add_styled_title(slide14, "Principaux Cas d'Usage")
    
    content14 = slide14.placeholders[1]
    use_cases = [
        "👤 Inscription et création de compte",
        "🔐 Connexion et authentification",
        "📄 Demande de document administratif",
        "🏗️ Soumission demande d'autorisation",
        "🚨 Signalement de problème urbain",
        "📰 Consultation des actualités",
        "📊 Suivi du statut des demandes",
        "👨‍💼 Gestion administrative (back-office)"
    ]
    add_bullet_points(content14, use_cases)
    print("✅ Slide 14: Cas d'usage")
    
    # ============== SECTION IV: DÉVELOPPEMENT ==============
    add_section_header(prs, "IV. DÉVELOPPEMENT", "Implémentation Technique de la Solution")
    print("✅ Section IV: Header")
    
    # SLIDE 15: Architecture Frontend
    slide15 = prs.slides.add_slide(prs.slide_layouts[1])
    add_styled_title(slide15, "Développement Frontend - React.js")
    
    content15 = slide15.placeholders[1]
    frontend_points = [
        "⚛️ React 18: Hooks et composants fonctionnels",
        "🎨 Tailwind CSS: Styling moderne et responsive",
        "🧭 React Router: Navigation côté client",
        "🔄 Context API: Gestion d'état globale",
        "📡 Axios: Communication API HTTP",
        "🔔 React-Toastify: Notifications utilisateur",
        "📱 Responsive Design: Mobile et desktop",
        "🚀 Optimisation: Lazy loading et code splitting"
    ]
    add_bullet_points(content15, frontend_points)
    print("✅ Slide 15: Frontend")
    
    # SLIDE 16: Architecture Backend
    slide16 = prs.slides.add_slide(prs.slide_layouts[1])
    add_styled_title(slide16, "Développement Backend - Node.js")
    
    content16 = slide16.placeholders[1]
    backend_points = [
        "🚀 Node.js + Express.js: Serveur rapide",
        "🗄️ MongoDB + Mongoose: ODM efficace",
        "🔐 JWT: Authentification sans session",
        "🛡️ bcrypt: Hachage sécurisé des mots de passe",
        "📧 Nodemailer: Service d'email automatisé",
        "📁 Multer: Gestion upload de fichiers",
        "🔒 CORS: Sécurité cross-origin",
        "⚡ Middleware: Validation et logging"
    ]
    add_bullet_points(content16, backend_points)
    print("✅ Slide 16: Backend")
    
    # SLIDE 17: Sécurité Implémentée
    slide17 = prs.slides.add_slide(prs.slide_layouts[1])
    add_styled_title(slide17, "Mesures de Sécurité Implémentées")
    
    content17 = slide17.placeholders[1]
    security_points = [
        "🔐 Authentification JWT sécurisée",
        "🛡️ Hachage bcrypt (12 rounds minimum)",
        "🌐 HTTPS: Chiffrement des communications",
        "🔒 Validation stricte des entrées utilisateur",
        "🚫 Protection CSRF et XSS",
        "⏱️ Rate limiting: Protection DDoS basique",
        "📝 Logging détaillé des actions",
        "🔑 Variables d'environnement sécurisées"
    ]
    add_bullet_points(content17, security_points)
    print("✅ Slide 17: Sécurité")
    
    # SLIDE 18: Fonctionnalités Clés Développées
    slide18 = prs.slides.add_slide(prs.slide_layouts[1])
    add_styled_title(slide18, "Fonctionnalités Clés Développées")
    
    content18 = slide18.placeholders[1]
    features_points = [
        "📝 Module Demandes Administratives",
        "🚨 Système de Signalements Citoyens",
        "📰 Gestionnaire d'Actualités Municipales",
        "👤 Profils Utilisateurs Personnalisés",
        "📊 Dashboard Administratif Complet",
        "🔍 Moteur de Recherche Intégré",
        "📧 Notifications Email Automatiques",
        "📱 Interface Mobile Responsive"
    ]
    add_bullet_points(content18, features_points)
    print("✅ Slide 18: Fonctionnalités")
    
    # ============== SPECIAL SLIDE: DEMONSTRATION ==============
    slide19 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout for video
    
    # Centered title at the top
    title_shape = slide19.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_shape.text_frame
    title_frame.text = "🎥 LIVE DEMONSTRATION"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(220, 20, 60)  # Bright red
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Main area for video
    video_placeholder = slide19.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(4))
    video_frame = video_placeholder.text_frame
    video_frame.text = """🎬 INSERT YOUR DEMONSTRATION VIDEO HERE
    
    💡 Demo Tips:
    • Recommended duration: 3-5 minutes
    • Show complete user journey
    • Demo both citizen AND admin sides
    • Highlight the modern interface
    • Present key features
    
    📱 Scenarios to demonstrate:
    → Registration and login
    → Document request
    → Problem reporting
    → Request status tracking
    → Administrator interface"""
    
    video_frame.paragraphs[0].font.size = Pt(20)
    video_frame.paragraphs[0].font.bold = True
    video_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    for i, p in enumerate(video_frame.paragraphs):
        if i > 0:
            p.font.size = Pt(14)
            p.alignment = PP_ALIGN.LEFT
    
    # Bordure décorative
    border_shape = slide19.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.4), Inches(1.9), Inches(7.2), Inches(4.2)
    )
    border_shape.fill.background()
    border_shape.line.color.rgb = RGBColor(220, 20, 60)
    border_shape.line.width = Pt(3)
    
    print("✅ Slide 19: 🎥 DÉMONSTRATION VIDÉO (SLIDE SPÉCIALE)")
    
    # ============== SECTION V: TESTS ==============
    add_section_header(prs, "V. TESTS ET VALIDATION", "Qualité et Fiabilité du Système")
    print("✅ Section V: Header")
    
    # SLIDE 20: Stratégie de Tests
    slide20 = prs.slides.add_slide(prs.slide_layouts[1])
    add_styled_title(slide20, "Stratégie de Tests Adoptée")
    
    content20 = slide20.placeholders[1]
    testing_points = [
        "🧪 Tests Unitaires: Fonctions critiques",
        "🔗 Tests d'Intégration: API endpoints",
        "👤 Tests Utilisateur: Parcours complets",
        "📱 Tests Responsive: Tous appareils",
        "🔒 Tests de Sécurité: Vulnérabilités",
        "⚡ Tests de Performance: Charge système",
        "♿ Tests d'Accessibilité: Standards WCAG",
        "🌐 Tests Cross-Browser: Compatibilité"
    ]
    add_bullet_points(content20, testing_points)
    print("✅ Slide 20: Tests")
    
    # SLIDE 21: Résultats des Tests
    slide21 = prs.slides.add_slide(prs.slide_layouts[1])
    add_styled_title(slide21, "Résultats des Tests et Métriques")
    
    content21 = slide21.placeholders[1]
    results_points = [
        "✅ 95% de couverture de code (tests unitaires)",
        "⚡ Temps de réponse moyen: 1.2 secondes",
        "👥 Support testé: 500 utilisateurs simultanés",
        "🔒 Aucune vulnérabilité critique détectée",
        "📱 Compatibilité: Chrome, Firefox, Safari, Edge",
        "♿ Score d'accessibilité: 92/100 (Lighthouse)",
        "🚀 Performance Lighthouse: 89/100",
        "✅ 0 erreurs de validation HTML/CSS"
    ]
    add_bullet_points(content21, results_points)
    print("✅ Slide 21: Résultats tests")
    
    # ============== SECTION VI: DÉPLOIEMENT ==============
    add_section_header(prs, "VI. DÉPLOIEMENT", "Mise en Production et Maintenance")
    print("✅ Section VI: Header")
    
    # SLIDE 22: Architecture de Déploiement
    slide22 = prs.slides.add_slide(prs.slide_layouts[1])
    add_styled_title(slide22, "Architecture de Déploiement")
    
    content22 = slide22.placeholders[1]
    deployment_points = [
        "☁️ Frontend: Vercel (CDN global)",
        "🚀 Backend: Vercel Serverless Functions",
        "🗄️ Base de données: MongoDB Atlas",
        "📧 Email: Service SMTP configuré",
        "🔒 SSL: Certificats automatiques",
        "🌐 DNS: Configuration domaine personnalisé",
        "📊 Monitoring: Logs et analytics intégrés",
        "🔄 CI/CD: Déploiement automatique Git"
    ]
    add_bullet_points(content22, deployment_points)
    print("✅ Slide 22: Déploiement")
    
    # SLIDE 23: Plan de Maintenance
    slide23 = prs.slides.add_slide(prs.slide_layouts[1])
    add_styled_title(slide23, "Plan de Maintenance et Support")
    
    content23 = slide23.placeholders[1]
    maintenance_points = [
        "🔄 Mises à jour sécuritaires mensuelles",
        "📊 Monitoring 24/7 des performances",
        "💾 Sauvegardes automatiques quotidiennes",
        "🛠️ Maintenance préventive trimestrielle",
        "📞 Support technique dédié",
        "📈 Analyses d'usage et optimisations",
        "🔐 Audits de sécurité semestriels",
        "📚 Documentation technique maintenue"
    ]
    add_bullet_points(content23, maintenance_points)
    print("✅ Slide 23: Maintenance")
    
    # ============== SECTION VII: IMPACTS ==============
    add_section_header(prs, "VII. IMPACTS ET BÉNÉFICES", "Valeur Ajoutée du Projet")
    print("✅ Section VII: Header")
    
    # SLIDE 24: Bénéfices pour les Citoyens
    slide24 = prs.slides.add_slide(prs.slide_layouts[1])
    add_styled_title(slide24, "Bénéfices pour les Citoyens")
    
    content24 = slide24.placeholders[1]
    citizen_benefits = [
        "🏠 Services accessibles 24h/24 depuis domicile",
        "⏱️ Réduction drastique des temps d'attente",
        "📱 Interface moderne et intuitive",
        "🔍 Suivi en temps réel des demandes",
        "💸 Économies de déplacements et frais",
        "📧 Notifications automatiques de statut",
        "🌐 Accessibilité pour personnes à mobilité réduite",
        "🤝 Communication améliorée avec la commune"
    ]
    add_bullet_points(content24, citizen_benefits)
    print("✅ Slide 24: Bénéfices citoyens")
    
    # SLIDE 25: Bénéfices pour l'Administration
    slide25 = prs.slides.add_slide(prs.slide_layouts[1])
    add_styled_title(slide25, "Bénéfices pour l'Administration")
    
    content25 = slide25.placeholders[1]
    admin_benefits = [
        "📊 Digitalisation complète des processus",
        "⚡ Traitement automatisé des demandes simples",
        "📈 Tableaux de bord analytiques détaillés",
        "💾 Archivage électronique sécurisé",
        "👥 Réduction de la charge d'accueil physique",
        "🔄 Workflows optimisés et standardisés",
        "📧 Communication automatisée avec citoyens",
        "💰 Réduction des coûts opérationnels"
    ]
    add_bullet_points(content25, admin_benefits)
    print("✅ Slide 25: Bénéfices administration")
    
    # SLIDE 26: Impact Économique et Social
    slide26 = prs.slides.add_slide(prs.slide_layouts[1])
    add_styled_title(slide26, "Impact Économique et Social")
    
    content26 = slide26.placeholders[1]
    economic_impact = [
        "💰 Réduction des coûts administratifs de 40%",
        "⏰ Gain de temps moyen: 2h par démarche",
        "🌱 Réduction de l'empreinte carbone (moins de déplacements)",
        "📋 Diminution de 80% de l'usage papier",
        "👥 Amélioration de la satisfaction citoyenne",
        "🎯 Modernisation de l'image municipale",
        "🤝 Renforcement de la confiance publique",
        "🏆 Positionnement d'Azrou comme ville innovante"
    ]
    add_bullet_points(content26, economic_impact)
    print("✅ Slide 26: Impact économique")
    
    # ============== SECTION VIII: PERSPECTIVES ==============
    add_section_header(prs, "VIII. PERSPECTIVES D'ÉVOLUTION", "Roadmap et Vision Future")
    print("✅ Section VIII: Header")
    
    # SLIDE 27: Évolutions à Court Terme
    slide27 = prs.slides.add_slide(prs.slide_layouts[1])
    add_styled_title(slide27, "Évolutions Prévues à Court Terme (6 mois)")
    
    content27 = slide27.placeholders[1]
    short_term = [
        "📱 Application mobile native (iOS/Android)",
        "💳 Paiement en ligne des taxes et redevances",
        "🔔 Notifications push en temps réel",
        "🗂️ Module de gestion électronique des documents",
        "📊 Dashboard avancé avec KPI détaillés",
        "🤖 Chatbot IA pour support automatisé",
        "🌐 Intégration API gouvernementales marocaines",
        "📧 Newsletter et communications ciblées"
    ]
    add_bullet_points(content27, short_term)
    print("✅ Slide 27: Court terme")
    
    # SLIDE 28: Vision à Long Terme
    slide28 = prs.slides.add_slide(prs.slide_layouts[1])
    add_styled_title(slide28, "Vision Smart City (2-5 ans)")
    
    content28 = slide28.placeholders[1]
    long_term = [
        "🏙️ Plateforme Smart City complète",
        "📡 IoT: Capteurs urbains connectés",
        "🚦 Gestion intelligente du trafic",
        "♻️ Optimisation de la gestion des déchets",
        "💡 Éclairage public intelligent",
        "📊 Big Data et analytics prédictives",
        "🤖 Intelligence artificielle intégrée",
        "🌍 Modèle réplicable pour autres communes"
    ]
    add_bullet_points(content28, long_term)
    print("✅ Slide 28: Long terme")
    
    # SLIDE 29: Défis et Recommandations
    slide29 = prs.slides.add_slide(prs.slide_layouts[1])
    add_styled_title(slide29, "Défis et Recommandations")
    
    content29 = slide29.placeholders[1]
    challenges = [
        "📚 Formation des agents municipaux",
        "👥 Accompagnement au changement des citoyens",
        "🔒 Renforcement continu de la cybersécurité",
        "💰 Budget dédié à l'évolution technologique",
        "🤝 Partenariats avec l'écosystème tech marocain",
        "📊 Mesure continue de la satisfaction utilisateur",
        "🔄 Mise à jour régulière des technologies",
        "🎯 Alignement avec la stratégie Maroc Digital 2030"
    ]
    add_bullet_points(content29, challenges)
    print("✅ Slide 29: Défis")
    
    # ============== SECTION IX: CONCLUSION ==============
    add_section_header(prs, "IX. CONCLUSION", "Synthèse et Ouverture")
    print("✅ Section IX: Header")
    
    # SLIDE 30: Synthèse du Projet
    slide30 = prs.slides.add_slide(prs.slide_layouts[1])
    add_styled_title(slide30, "Synthèse du Projet")
    
    content30 = slide30.placeholders[1]
    synthesis = [
        "✅ Objectifs pleinement atteints",
        "🚀 Solution moderne et évolutive développée",
        "💡 Technologies actuelles maîtrisées",
        "🎯 Besoins utilisateurs satisfaits",
        "📊 Résultats mesurables et concrets",
        "🤝 Impact positif sur citoyens et administration",
        "🌱 Bases solides pour évolutions futures",
        "🏆 Contribution à la transformation numérique"
    ]
    add_bullet_points(content30, synthesis)
    print("✅ Slide 30: Synthèse")
    
    # SLIDE 31: Apports Personnels et Apprentissages
    slide31 = prs.slides.add_slide(prs.slide_layouts[1])
    add_styled_title(slide31, "Apports Personnels et Apprentissages")
    
    content31 = slide31.placeholders[1]
    learning = [
        "💻 Maîtrise du stack MERN complet",
        "🏛️ Compréhension des enjeux municipaux",
        "👥 Gestion de projet en autonomie",
        "🔍 Analyse des besoins utilisateurs",
        "🎨 Design UX/UI orienté utilisabilité",
        "🔒 Implémentation de la sécurité web",
        "☁️ Déploiement et DevOps cloud",
        "🤝 Communication avec parties prenantes"
    ]
    add_bullet_points(content31, learning)
    print("✅ Slide 31: Apprentissages")
    
    # SLIDE 32: Questions et Discussion
    slide32 = prs.slides.add_slide(prs.slide_layouts[0])  # Title layout pour finaliser
    title32 = slide32.shapes.title
    subtitle32 = slide32.placeholders[1]
    
    title32.text = "MERCI POUR VOTRE ATTENTION"
    subtitle32.text = """Questions et Discussion

🤔 Vos Questions
💭 Retours et Suggestions
🔍 Approfondissements Techniques

Anas BENOMAR
Développeur Full-Stack
📧 Contact pour démonstration approfondie"""
    
    # Style final
    title32.text_frame.paragraphs[0].font.size = Pt(36)
    title32.text_frame.paragraphs[0].font.bold = True
    title32.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    title32.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    for p in subtitle32.text_frame.paragraphs:
        p.font.size = Pt(18)
        p.alignment = PP_ALIGN.CENTER
    
    print("✅ Slide 32: Questions et discussion")
    
    # Save the presentation
    filename = "Presentation_PFE_Extended_English_Azrou_Anas_BENOMAR.pptx"
    prs.save(filename)
    
    print("=" * 60)
    print(f"🎉 PRESENTATION GENERATED SUCCESSFULLY!")
    print(f"📁 File: {filename}")
    print(f"📊 Number of slides: {len(prs.slides)}")
    print(f"⏱️ Estimated presentation duration: 25-30 minutes")
    print("=" * 60)
    
    return filename

if __name__ == "__main__":
    try:
        create_comprehensive_presentation()
        print("\n✅ Generation completed successfully!")
        print("💡 Tips:")
        print("   • Add your application screenshots")
        print("   • Customize colors according to your brand")
        print("   • Practice your presentation several times")
        print("   • Prepare for technical questions")
        print("   🎥 Don't forget to insert your demo video!")
        
    except Exception as e:
        print(f"❌ Error during generation: {e}")
        import traceback
        traceback.print_exc()
